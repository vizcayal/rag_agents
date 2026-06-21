import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_deck():
    prs = Presentation()
    
    # Configure 16:9 widescreen dimensions (standard for LinkedIn and PPTX)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Theme color definitions
    bg_color = RGBColor(13, 17, 23)        # Deep slate background
    title_color = RGBColor(88, 166, 255)   # Cyan/Blue accent
    text_color = RGBColor(201, 209, 217)   # Soft white/gray body text
    sub_color = RGBColor(139, 148, 158)    # Muted gray text
    highlight_color = RGBColor(188, 140, 255) # Purple accent
    
    # Slide 1: Cover
    slide_layout = prs.slide_layouts[6] # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    
    # Cover text box (Centered)
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Cómo Desplegar un Sistema RAG\nen Producción en AWS 🚀"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Segoe UI'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.space_after = Pt(20)
    
    p2 = tf.add_paragraph()
    p2.text = "Arquitectura Serverless y Contenedores Autogestionados (Frontend + Backend)"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = 'Segoe UI'
    p2.font.size = Pt(20)
    p2.font.color.rgb = text_color
    
    # Helper to add standard content slides
    def add_content_slide(title_text, bullets, subtitle=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color
        
        # Title Box
        t_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1.2))
        tf = t_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Segoe UI'
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = title_color
        
        if subtitle:
            p.space_after = Pt(4)
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.name = 'Segoe UI'
            p2.font.size = Pt(16)
            p2.font.color.rgb = sub_color
            
        # Bullets / Body text
        c_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(11.833), Inches(4.5))
        cf = c_box.text_frame
        cf.word_wrap = True
        
        for idx, bullet in enumerate(bullets):
            p_b = cf.paragraphs[0] if idx == 0 else cf.add_paragraph()
            p_b.text = bullet
            p_b.font.name = 'Segoe UI'
            p_b.font.size = Pt(18)
            p_b.font.color.rgb = text_color
            p_b.space_before = Pt(12)
            p_b.space_after = Pt(6)
            
            # Simple custom bullet indent
            p_b.level = 0
            
        return slide

    # Slide 2: El Reto
    add_content_slide(
        title_text="El Reto: Pasar del Entorno Local a la Nube",
        subtitle="Por qué los prototipos locales fallan en la vida real",
        bullets=[
            "🛑 Escalabilidad: Las consultas simultáneas saturan los recursos locales de CPU/GPU.",
            "🛑 Seguridad y Acceso: Falta de aislamiento de red y de control de identidades (IAM).",
            "🛑 Actualizaciones manuales: Cada cambio de código interrumpe el servicio.",
            "💡 Solución: Una arquitectura híbrida en AWS combinando Serverless para la base vectorial y Fargate para la interfaz de usuario."
        ]
    )
    
    # Slide 3: Arquitectura General
    add_content_slide(
        title_text="Arquitectura General en AWS",
        subtitle="Cómo se conectan e interactúan las piezas del sistema",
        bullets=[
            "🖥️ Frontend: Streamlit (Python) ejecutado en contenedores de AWS ECS Fargate.",
            "⚙️ Orquestador Backend: Agente LangGraph integrado con APIs en Bedrock.",
            "🔍 Vector Store: Base de datos vectorial en Amazon OpenSearch Serverless (AOSS).",
            "📄 Ingesta de Datos: PDFs almacenados en S3 sincronizados con Bedrock Knowledge Base.",
            "🔐 Redes y Coseguro: Balanceador ALB expone el Frontend y políticas IAM gestionan los accesos."
        ]
    )

    # Slide 4: El Frontend (ECS Fargate)
    add_content_slide(
        title_text="Frontend: Contenedorizando con Docker y ECS",
        subtitle="Alojamiento resiliente de la UI a mínimo coste",
        bullets=[
            "🐳 Contenedorizado: Imagen ligera con Python 3.11-slim y comprobación de salud activa.",
            "🌐 Redes de Bajo Coste: VPC sin NAT Gateways (tareas ECS en subnets públicas con IP pública).",
            "⚖️ Balanceador de Carga (ALB): Rutea tráfico HTTP y distribuye peticiones entre contenedores.",
            "🔄 Auto-Healing: Si la UI falla, ECS Fargate recrea la tarea en segundos."
        ]
    )

    # Slide 5: El Backend (LangGraph + OpenSearch Serverless)
    add_content_slide(
        title_text="Backend: Base Vectorial y Razonamiento RAG",
        subtitle="La base de conocimiento sin gestión de servidores",
        bullets=[
            "🔍 OpenSearch Serverless: Colección vectorial con indexación K.N.N y embeddings de Amazon Titan.",
            "📚 Ingesta Inteligente: Conexión nativa de Bedrock KB con S3 para chunking y guardado automático.",
            "🧠 Inferencia RAG: Modelo Amazon Nova Lite para sintetizar respuestas.",
            "🛡️ Seguridad Estricta: Permisos de invocación cifrados mediante IAM Policies de mínimos privilegios."
        ]
    )

    # Slide 6: Flujo de Comunicación paso a paso
    add_content_slide(
        title_text="El Flujo de una Consulta RAG",
        subtitle="Paso a paso desde que el usuario pregunta",
        bullets=[
            "1️⃣ El usuario escribe una pregunta en la interfaz de Streamlit.",
            "2️⃣ Streamlit invoca al agente de Bedrock/LangGraph a través de boto3.",
            "3️⃣ El agente recupera fragmentos indexados en la colección de OpenSearch Serverless.",
            "4️⃣ El modelo Nova Lite de Bedrock genera la respuesta basada únicamente en las fuentes.",
            "5️⃣ La respuesta y las citas textuales de origen se devuelven en streaming a la interfaz."
        ]
    )

    # Slide 7: Validación de Respuestas (Evitando Alucinaciones)
    add_content_slide(
        title_text="Validación Groundedness en LangGraph",
        subtitle="Garantizar que la IA no invente datos",
        bullets=[
            "✅ Verificación RAG: Filtros que analizan la respuesta frente al texto original del PDF.",
            "📊 Semáforo de Confianza:",
            "   🟢 Grounded: El texto recuperado sustenta 100% la respuesta.",
            "   🟡 Respuesta Parcial: Contiene inferencias no directas.",
            "   🔴 No Verificado: Alucinación potencial - la respuesta no se ancla en los documentos.",
            "🔗 Citación de Fuentes: Muestra el fragmento oficial con su origen en S3."
        ]
    )

    # Slide 8: Despliegue en 3 Pasos
    add_content_slide(
        title_text="Despliegue y CI/CD Automatizado",
        subtitle="Actualizaciones sencillas y continuas en caliente",
        bullets=[
            "🐳 Paso 1: Compilación local de la UI:\n   docker build -t ecr-repo-url:latest .",
            "🚀 Paso 2: Push al registro seguro de ECR:\n   docker push ecr-repo-url:latest",
            "⚡ Paso 3: Despliegue Zero-Downtime en ECS Fargate:\n   aws ecs update-service --cluster UI-Cluster --service Streamlit-Service --force-new-deployment"
        ]
    )

    # Slide 9: Conclusiones y Llamado a la Acción (CTA)
    add_content_slide(
        title_text="Lecciones del Proyecto y CTA",
        subtitle="Claves para llevar a casa",
        bullets=[
            "💡 Serverless es Rey: AOSS y Bedrock eliminan la sobrecarga de administración.",
            "🛡️ La confianza se mide: La validación RAG es obligatoria para sistemas profesionales.",
            "💬 ¿Cómo manejas las alucinaciones en tus aplicaciones de IA?",
            "👇 ¡Escribe en los comentarios o comparte para ayudar a otros desarrolladores!"
        ]
    )
    
    # Save presentation
    filename = "Despliegue_RAG_AWS.pptx"
    prs.save(filename)
    print(f"Presentation saved successfully as '{filename}'")

if __name__ == "__main__":
    create_deck()
